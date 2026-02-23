#!/usr/bin/env python3
"""
Generate Valley of Fire Project Roadmap PowerPoint for OP2 and OP3.
Based on PWS deliverables. Timeline: OP2 (Feb-May 2026), OP3 (May-Aug 2026).

Usage:
    python3 vof/generate_vof_roadmap.py                # dark (default)
    python3 vof/generate_vof_roadmap.py --theme light
    python3 vof/generate_vof_roadmap.py --theme dark
"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Timeline
# ---------------------------------------------------------------------------
QUARTERS = ["Q1 2026", "Q2 2026", "Q3 2026"]
MONTHS = ["Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"]

# ---------------------------------------------------------------------------
# Layout
# ---------------------------------------------------------------------------
TABLE_LEFT = 0.5
TABLE_TOP = 1.08
PHASE_COL_WIDTH = 2.7
MONTH_COLS = len(MONTHS)
CHART_WIDTH = 5.6
MONTH_WIDTH = CHART_WIDTH / MONTH_COLS
CHART_LEFT = TABLE_LEFT + PHASE_COL_WIDTH
DUE_DATE_COL_WIDTH = 1.1
HEADER_HEIGHT = 0.4
ROW_HEIGHT = 0.38
BAR_HEIGHT = 0.22
BAR_PADDING = 0.03

# ---------------------------------------------------------------------------
# Themes
# ---------------------------------------------------------------------------
THEMES = {
    "dark": {
        "slide_bg": RGBColor(15, 23, 42),
        "header_bg": RGBColor(30, 41, 59),
        "header_text": RGBColor(255, 255, 255),
        "title_text": RGBColor(226, 232, 240),
        "subtitle_text": RGBColor(148, 163, 184),
        "period_text": RGBColor(148, 163, 184),
        "row_even": RGBColor(30, 41, 59),
        "row_odd": RGBColor(51, 65, 85),
        "task_text": RGBColor(226, 232, 240),
        "due_text": RGBColor(148, 163, 184),
        "footer_text": RGBColor(100, 116, 139),
        "bar_colors": {
            "App Maturity": RGBColor(34, 197, 94),
            "Data & Integration": RGBColor(249, 115, 22),
            "Platform Sustainment": RGBColor(139, 92, 246),
            "Program Deliverables": RGBColor(236, 72, 153),
            "Optional / Growth": RGBColor(148, 163, 184),
        },
    },
    "light": {
        "slide_bg": RGBColor(255, 255, 255),
        "header_bg": RGBColor(241, 245, 249),
        "header_text": RGBColor(30, 41, 59),
        "title_text": RGBColor(15, 23, 42),
        "subtitle_text": RGBColor(71, 85, 105),
        "period_text": RGBColor(71, 85, 105),
        "row_even": RGBColor(255, 255, 255),
        "row_odd": RGBColor(248, 250, 252),
        "task_text": RGBColor(30, 41, 59),
        "due_text": RGBColor(71, 85, 105),
        "footer_text": RGBColor(148, 163, 184),
        "bar_colors": {
            "App Maturity": RGBColor(22, 163, 74),
            "Data & Integration": RGBColor(234, 88, 12),
            "Platform Sustainment": RGBColor(124, 58, 237),
            "Program Deliverables": RGBColor(219, 39, 119),
            "Optional / Growth": RGBColor(100, 116, 139),
        },
    },
}

# ---------------------------------------------------------------------------
# Tasks: (phase, task, start_m, end_m, is_milestone, due_date)
# ---------------------------------------------------------------------------
TASKS = [
    ("App Maturity", "IR Review & INTEL Agent v0→v1", 0, 1, False, "1 Mar 2026"),
    ("Data & Integration", "Complete ServiceNow data connection", 0, 3, False, "TBD"),
    ("Data & Integration", "New data connections (OP3 optional)", 4, 6, False, None),
    ("Platform Sustainment", "Platform licenses, updates & 9-app sustainment", 0, 6, False, None),
    ("Program Deliverables", "OP2 PoC Plan", 0, 0, True, "20 Feb 2026"),
    ("Program Deliverables", "OP2 Closeout & Dataset Transition", 2, 3, False, "14 May 2026"),
    ("Program Deliverables", "OP3 Closeout & Dataset Transition", 5, 6, False, "14 Aug 2026"),
    ("Program Deliverables", "Monthly MTR / MFR", 0, 6, False, "7 Feb 2026"),
    ("Optional / Growth", "Scoping for new applications", 4, 6, False, None),
    ("Optional / Growth", "Up to 3 new apps (CLIN 0013)", 5, 6, False, "14 Aug 2026"),
]


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def main(theme_name: str = "dark"):
    theme = THEMES[theme_name]
    bar_colors = theme["bar_colors"]
    fallback_bar = bar_colors["Optional / Growth"]

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["slide_bg"])

    # Subtitle
    subtitle = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.15), Inches(9), Inches(0.3)
    )
    sp = subtitle.text_frame.paragraphs[0]
    sp.text = "Valley of Fire · Project ASCEND · HC1084-25-0001"
    sp.font.size = Pt(9)
    sp.font.color.rgb = theme["subtitle_text"]

    # Title
    title = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.4), Inches(9), Inches(0.55)
    )
    tp = title.text_frame.paragraphs[0]
    tp.text = "Project Roadmap: Option Period 2 & 3"
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = theme["title_text"]

    # Period callout
    period_box = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.8), Inches(9), Inches(0.22)
    )
    pp = period_box.text_frame.paragraphs[0]
    pp.text = "OP2: 15 Feb – 14 May 2026  ·  OP3: 15 May – 14 Aug 2026"
    pp.font.size = Pt(8)
    pp.font.color.rgb = theme["period_text"]

    # Table
    total_rows = 2 + len(TASKS)
    total_cols = 1 + MONTH_COLS + 1
    due_col = total_cols - 1
    table_shape = slide.shapes.add_table(
        total_rows, total_cols,
        Inches(TABLE_LEFT), Inches(TABLE_TOP),
        Inches(PHASE_COL_WIDTH + CHART_WIDTH + DUE_DATE_COL_WIDTH),
        Inches(HEADER_HEIGHT * 2 + len(TASKS) * ROW_HEIGHT),
    )
    table = table_shape.table

    # Header rows
    for r in range(2):
        for c in range(total_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["header_bg"]
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = theme["header_text"]
            para.font.size = Pt(12)
            para.font.bold = (r == 0)
            para.alignment = PP_ALIGN.CENTER

    table.cell(0, 0).text = "Deliverable"
    table.cell(0, 1).text = QUARTERS[0]
    table.cell(0, 1).merge(table.cell(0, 2))
    table.cell(0, 3).text = QUARTERS[1]
    table.cell(0, 3).merge(table.cell(0, 4))
    table.cell(0, 5).text = QUARTERS[2]
    table.cell(0, 5).merge(table.cell(0, 7))
    table.cell(0, due_col).text = "Due Date"
    table.cell(0, due_col).merge(table.cell(1, due_col))

    table.cell(1, 0).text = ""
    for c, m in enumerate(MONTHS):
        table.cell(1, 1 + c).text = m

    table.columns[0].width = Inches(PHASE_COL_WIDTH)
    for c in range(1, due_col):
        table.columns[c].width = Inches(MONTH_WIDTH)
    table.columns[due_col].width = Inches(DUE_DATE_COL_WIDTH)

    # Data rows
    for row in range(2, total_rows):
        task = TASKS[row - 2]
        bg = theme["row_odd"] if (row - 2) % 2 == 1 else theme["row_even"]
        for c in range(total_cols):
            cell = table.cell(row, c)
            if c == 0:
                cell.text = task[1]
            elif c == due_col and task[5]:
                cell.text = task[5]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            if c == 0:
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.color.rgb = theme["task_text"]
                para.alignment = PP_ALIGN.LEFT
            elif c == due_col:
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.color.rgb = theme["due_text"]
                para.alignment = PP_ALIGN.RIGHT

    # Task bars
    chart_top = TABLE_TOP + HEADER_HEIGHT * 2
    for i, task in enumerate(TASKS):
        phase, _, start_m, end_m, is_milestone, _ = task
        color = bar_colors.get(phase, fallback_bar)
        row_top = chart_top + i * ROW_HEIGHT
        center_y = row_top + ROW_HEIGHT / 2

        if is_milestone:
            size = 0.18
            center_x = CHART_LEFT + (start_m + 0.5) * MONTH_WIDTH
            shape = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(center_x - size / 2), Inches(center_y - size / 2),
                Inches(size), Inches(size),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
        else:
            bar_top = row_top + (ROW_HEIGHT - BAR_HEIGHT) / 2
            bar_left = CHART_LEFT + start_m * MONTH_WIDTH + BAR_PADDING
            span = end_m - start_m + 1
            bar_width = max(0.25, span * MONTH_WIDTH - 2 * BAR_PADDING)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_left), Inches(bar_top),
                Inches(bar_width), Inches(BAR_HEIGHT),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

    # Footer
    footer = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(7.1), Inches(9), Inches(0.25)
    )
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Scale AI  ·  Source: PWS Mod P00006"
    fp.font.size = Pt(8)
    fp.font.color.rgb = theme["footer_text"]

    # Save
    out_dir = Path(__file__).resolve().parent / "output"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / f"VoF_Roadmap_OP2_OP3_{theme_name}.pptx"
    prs.save(str(out_path))
    print(f"Created {out_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate VoF OP2/OP3 roadmap slide")
    parser.add_argument(
        "--theme", choices=["dark", "light"], default="dark",
        help="Color theme (default: dark)",
    )
    args = parser.parse_args()
    main(args.theme)
