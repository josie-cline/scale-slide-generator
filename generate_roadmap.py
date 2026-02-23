#!/usr/bin/env python3
"""
Generate a Gantt-style roadmap PowerPoint slide for any Scale program.

Usage:
    python3 generate_roadmap.py                # dark theme (default)
    python3 generate_roadmap.py --theme light
    python3 generate_roadmap.py --theme dark

Customization:
    Edit PROGRAM, QUARTERS, MONTHS, TASKS, and PHASES below to match
    your program. Or ask Cursor to update them for you:

        "Update the roadmap for [program name]. Here are the tasks: ..."
"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# ═══════════════════════════════════════════════════════════════════════════
# PROGRAM DATA — Edit this section for your program
# ═══════════════════════════════════════════════════════════════════════════

PROGRAM = {
    "name": "Program Name",
    "subtitle": "Scale AI · Contract HC0000-00-0000",
    "title": "Project Roadmap: Option Period 2 & 3",
    "period": "OP2: 15 Feb – 14 May 2026  ·  OP3: 15 May – 14 Aug 2026",
    "footer": "Scale AI  ·  Source: PWS",
}

QUARTERS = ["Q1 2026", "Q2 2026", "Q3 2026"]
MONTHS = ["Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"]

# Phases — each phase gets a distinct color on the Gantt bars.
# Add or remove phases as needed. Keys here must match the phase
# names used in TASKS below.
PHASES = [
    "Development",
    "Integration",
    "Sustainment",
    "Deliverables",
    "Optional",
]

# Tasks: (phase, task_name, start_month_idx, end_month_idx, is_milestone, due_date)
#   - start/end month indices are 0-based into MONTHS above
#   - is_milestone: True renders a diamond instead of a bar
#   - due_date: string shown in the Due Date column, or None to leave blank
TASKS = [
    ("Development", "Application refinement v0 → v1", 0, 1, False, "1 Mar 2026"),
    ("Integration", "Complete data source connection", 0, 3, False, "TBD"),
    ("Integration", "New data connections (optional)", 4, 6, False, None),
    ("Sustainment", "Platform licenses, updates & sustainment", 0, 6, False, None),
    ("Deliverables", "PoC Plan", 0, 0, True, "20 Feb 2026"),
    ("Deliverables", "OP2 Closeout & Transition", 2, 3, False, "14 May 2026"),
    ("Deliverables", "OP3 Closeout & Transition", 5, 6, False, "14 Aug 2026"),
    ("Deliverables", "Monthly Status Reviews", 0, 6, False, None),
    ("Optional", "Scoping for new capabilities", 4, 6, False, None),
    ("Optional", "New capability development", 5, 6, False, "14 Aug 2026"),
]


# ═══════════════════════════════════════════════════════════════════════════
# LAYOUT — Adjust if you need a wider/taller slide
# ═══════════════════════════════════════════════════════════════════════════

TABLE_LEFT = 0.5
TABLE_TOP = 1.08
PHASE_COL_WIDTH = 2.7
CHART_WIDTH = 5.6
DUE_DATE_COL_WIDTH = 1.1
HEADER_HEIGHT = 0.4
ROW_HEIGHT = 0.38
BAR_HEIGHT = 0.22
BAR_PADDING = 0.03

MONTH_COLS = len(MONTHS)
MONTH_WIDTH = CHART_WIDTH / MONTH_COLS
CHART_LEFT = TABLE_LEFT + PHASE_COL_WIDTH


# ═══════════════════════════════════════════════════════════════════════════
# THEMES — Color palettes for dark and light mode
# ═══════════════════════════════════════════════════════════════════════════

# Five bar colors that cycle across phases. Add more if you have > 5 phases.
_BAR_PALETTE_DARK = [
    RGBColor(34, 197, 94),    # green
    RGBColor(249, 115, 22),   # orange
    RGBColor(139, 92, 246),   # violet
    RGBColor(236, 72, 153),   # magenta
    RGBColor(148, 163, 184),  # slate
]
_BAR_PALETTE_LIGHT = [
    RGBColor(22, 163, 74),
    RGBColor(234, 88, 12),
    RGBColor(124, 58, 237),
    RGBColor(219, 39, 119),
    RGBColor(100, 116, 139),
]


def _build_bar_colors(phases, palette):
    return {phase: palette[i % len(palette)] for i, phase in enumerate(phases)}


THEMES = {
    "dark": {
        "slide_bg": RGBColor(15, 23, 42),
        "header_bg": RGBColor(30, 41, 59),
        "header_text": RGBColor(255, 255, 255),
        "title_text": RGBColor(226, 232, 240),
        "subtitle_text": RGBColor(148, 163, 184),
        "period_text": RGBColor(148, 163, 184),
        "row_even": RGBColor(30, 41, 59),
        "row_odd": RGBColor(51, 65, 85),
        "task_text": RGBColor(226, 232, 240),
        "due_text": RGBColor(148, 163, 184),
        "footer_text": RGBColor(100, 116, 139),
    },
    "light": {
        "slide_bg": RGBColor(255, 255, 255),
        "header_bg": RGBColor(241, 245, 249),
        "header_text": RGBColor(30, 41, 59),
        "title_text": RGBColor(15, 23, 42),
        "subtitle_text": RGBColor(71, 85, 105),
        "period_text": RGBColor(71, 85, 105),
        "row_even": RGBColor(255, 255, 255),
        "row_odd": RGBColor(248, 250, 252),
        "task_text": RGBColor(30, 41, 59),
        "due_text": RGBColor(71, 85, 105),
        "footer_text": RGBColor(148, 163, 184),
    },
}


# ═══════════════════════════════════════════════════════════════════════════
# GENERATOR — You shouldn't need to edit below this line
# ═══════════════════════════════════════════════════════════════════════════

def _set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def main(theme_name: str = "dark"):
    theme = THEMES[theme_name]
    palette = _BAR_PALETTE_DARK if theme_name == "dark" else _BAR_PALETTE_LIGHT
    bar_colors = _build_bar_colors(PHASES, palette)
    fallback = RGBColor(148, 163, 184)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["slide_bg"])

    # Subtitle
    sub = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.15), Inches(9), Inches(0.3)
    )
    p = sub.text_frame.paragraphs[0]
    p.text = PROGRAM["subtitle"]
    p.font.size = Pt(9)
    p.font.color.rgb = theme["subtitle_text"]

    # Title
    ttl = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.4), Inches(9), Inches(0.55)
    )
    p = ttl.text_frame.paragraphs[0]
    p.text = PROGRAM["title"]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = theme["title_text"]

    # Period callout
    per = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.8), Inches(9), Inches(0.22)
    )
    p = per.text_frame.paragraphs[0]
    p.text = PROGRAM["period"]
    p.font.size = Pt(8)
    p.font.color.rgb = theme["period_text"]

    # Table
    total_rows = 2 + len(TASKS)
    total_cols = 1 + MONTH_COLS + 1
    due_col = total_cols - 1

    table_shape = slide.shapes.add_table(
        total_rows, total_cols,
        Inches(TABLE_LEFT), Inches(TABLE_TOP),
        Inches(PHASE_COL_WIDTH + CHART_WIDTH + DUE_DATE_COL_WIDTH),
        Inches(HEADER_HEIGHT * 2 + len(TASKS) * ROW_HEIGHT),
    )
    table = table_shape.table

    # Header rows
    for r in range(2):
        for c in range(total_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["header_bg"]
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = theme["header_text"]
            para.font.size = Pt(12)
            para.font.bold = (r == 0)
            para.alignment = PP_ALIGN.CENTER

    # Row 0: Quarters (auto-merge based on count)
    table.cell(0, 0).text = "Deliverable"
    q_start = 1
    months_per_q = MONTH_COLS // len(QUARTERS) if QUARTERS else MONTH_COLS
    for i, q in enumerate(QUARTERS):
        c_start = 1 + i * months_per_q
        c_end = min(1 + (i + 1) * months_per_q - 1, MONTH_COLS)
        table.cell(0, c_start).text = q
        if c_end > c_start:
            table.cell(0, c_start).merge(table.cell(0, c_end))

    table.cell(0, due_col).text = "Due Date"
    table.cell(0, due_col).merge(table.cell(1, due_col))

    # Row 1: Months
    table.cell(1, 0).text = ""
    for c, m in enumerate(MONTHS):
        table.cell(1, 1 + c).text = m

    # Column widths
    table.columns[0].width = Inches(PHASE_COL_WIDTH)
    for c in range(1, due_col):
        table.columns[c].width = Inches(MONTH_WIDTH)
    table.columns[due_col].width = Inches(DUE_DATE_COL_WIDTH)

    # Data rows
    for row in range(2, total_rows):
        task = TASKS[row - 2]
        bg = theme["row_odd"] if (row - 2) % 2 == 1 else theme["row_even"]
        for c in range(total_cols):
            cell = table.cell(row, c)
            if c == 0:
                cell.text = task[1]
            elif c == due_col and task[5]:
                cell.text = task[5]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            if c == 0:
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.color.rgb = theme["task_text"]
                para.alignment = PP_ALIGN.LEFT
            elif c == due_col:
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.color.rgb = theme["due_text"]
                para.alignment = PP_ALIGN.RIGHT

    # Task bars
    chart_top = TABLE_TOP + HEADER_HEIGHT * 2
    for i, task in enumerate(TASKS):
        phase, _, start_m, end_m, is_milestone, _ = task
        color = bar_colors.get(phase, fallback)
        row_top = chart_top + i * ROW_HEIGHT
        center_y = row_top + ROW_HEIGHT / 2

        if is_milestone:
            size = 0.18
            center_x = CHART_LEFT + (start_m + 0.5) * MONTH_WIDTH
            shape = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(center_x - size / 2), Inches(center_y - size / 2),
                Inches(size), Inches(size),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
        else:
            bar_top = row_top + (ROW_HEIGHT - BAR_HEIGHT) / 2
            bar_left = CHART_LEFT + start_m * MONTH_WIDTH + BAR_PADDING
            span = end_m - start_m + 1
            bar_width = max(0.25, span * MONTH_WIDTH - 2 * BAR_PADDING)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_left), Inches(bar_top),
                Inches(bar_width), Inches(BAR_HEIGHT),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

    # Footer
    ftr = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(7.1), Inches(9), Inches(0.25)
    )
    p = ftr.text_frame.paragraphs[0]
    p.text = PROGRAM["footer"]
    p.font.size = Pt(8)
    p.font.color.rgb = theme["footer_text"]

    # Save
    out_dir = Path(__file__).resolve().parent / "output"
    out_dir.mkdir(exist_ok=True)
    safe_name = PROGRAM["name"].replace(" ", "_")
    out_path = out_dir / f"{safe_name}_Roadmap_{theme_name}.pptx"
    prs.save(str(out_path))
    print(f"Created {out_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate a Gantt-style roadmap slide for any Scale program"
    )
    parser.add_argument(
        "--theme", choices=["dark", "light"], default="dark",
        help="Color theme (default: dark)",
    )
    args = parser.parse_args()
    main(args.theme)
