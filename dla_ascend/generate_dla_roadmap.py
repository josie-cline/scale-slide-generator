"""
Generate DLA Roadmap PowerPoint for Option Period 2 & 3.
Recreates the Gantt chart slide from the PWS.

Usage:
    python3 dla_ascend/generate_dla_roadmap.py                # dark (default)
    python3 dla_ascend/generate_dla_roadmap.py --theme light
    python3 dla_ascend/generate_dla_roadmap.py --theme dark
"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Layout
# ---------------------------------------------------------------------------
SLIDE_WIDTH = 10
TABLE_LEFT = 0.5
TABLE_TOP = 1.0
PHASE_COL_WIDTH = 2.0
MONTH_COLS = 12
CHART_WIDTH = 7.0
MONTH_WIDTH = CHART_WIDTH / MONTH_COLS
CHART_LEFT = TABLE_LEFT + PHASE_COL_WIDTH
HEADER_HEIGHT = 0.35
ROW_HEIGHT = 0.3
BAR_HEIGHT = 0.22

# ---------------------------------------------------------------------------
# Timeline
# ---------------------------------------------------------------------------
QUARTERS = ["Q2 2025", "Q3 2025", "Q4 2025", "Q1 2026"]
MONTHS = ["APR", "MAY", "JUN", "JUL", "AUG", "SEP",
          "OCT", "NOV", "DEC", "JAN", "FEB", "MAR"]

# ---------------------------------------------------------------------------
# Themes
# ---------------------------------------------------------------------------
THEMES = {
    "dark": {
        "slide_bg": RGBColor(15, 23, 42),
        "header_bg": RGBColor(30, 41, 59),
        "header_text": RGBColor(255, 255, 255),
        "phase_bg": RGBColor(30, 41, 59),
        "phase_text": RGBColor(255, 255, 255),
        "title_text": RGBColor(226, 232, 240),
        "row_even": RGBColor(30, 41, 59),
        "row_odd": RGBColor(51, 65, 85),
        "bar_text": RGBColor(226, 232, 240),
        "bar_colors": {
            "Phase I": RGBColor(139, 92, 246),
            "Phase II": RGBColor(56, 189, 248),
            "Phase III": RGBColor(251, 113, 133),
            "Other": RGBColor(148, 163, 184),
        },
    },
    "light": {
        "slide_bg": RGBColor(255, 255, 255),
        "header_bg": RGBColor(241, 245, 249),
        "header_text": RGBColor(30, 41, 59),
        "phase_bg": RGBColor(51, 65, 85),
        "phase_text": RGBColor(255, 255, 255),
        "title_text": RGBColor(15, 23, 42),
        "row_even": RGBColor(255, 255, 255),
        "row_odd": RGBColor(248, 250, 252),
        "bar_text": RGBColor(30, 41, 59),
        "bar_colors": {
            "Phase I": RGBColor(180, 160, 220),
            "Phase II": RGBColor(170, 200, 230),
            "Phase III": RGBColor(230, 170, 180),
            "Other": RGBColor(160, 160, 160),
        },
    },
}

# ---------------------------------------------------------------------------
# Tasks: (phase_name, task_label, start_month_idx, end_month_idx)
# ---------------------------------------------------------------------------
TASKS = [
    ("Phase I: Scale Gov Onboarding", "Engage DLA Stakeholders on use cases", 0, 2),
    ("Phase I: Scale Gov Onboarding", "Ingest data in Scale Gov / Create API connections", 1, 4),
    ("Phase I: Scale Gov Onboarding", "Generate SITREPS", 2, 4),
    ("Phase I: Scale Gov Onboarding", "Generate Charts", 2, 5),
    ("Phase II: DLA Cloud Deployment", "Deploy Donovan to DLA Cloud", 1, 4),
    ("Phase II: DLA Cloud Deployment", "Support DLA ATO Upgrade", 5, 7),
    ("Phase III: Cloud Data Connections and Reports Gen", "Engage DLA Stakeholders on use cases", 6, 11),
    ("Phase III: Cloud Data Connections and Reports Gen", "Ingest data in Scale Gov / Create API connections", 7, 10),
    ("Phase III: Cloud Data Connections and Reports Gen", "Generate SITREPS", 7, 9),
    ("Phase III: Cloud Data Connections and Reports Gen", "Generate Charts", 8, 10),
    ("Phase III: Cloud Data Connections and Reports Gen", "Measure performance increase with Donovan", 8, 10),
    ("Other", "Monthly MTR / MFRs", 0, 11),
    ("Other", "ID other workflows that could be improved through LLM integration / overdeliver", 1, 11),
]


def _phase_key(phase_name: str) -> str:
    for key in ("Phase I", "Phase II", "Phase III"):
        if key in phase_name:
            return key
    return "Other"


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def main(theme_name: str = "dark"):
    theme = THEMES[theme_name]
    bar_colors = theme["bar_colors"]

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["slide_bg"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(TABLE_LEFT), Inches(0.25), Inches(9), Inches(0.5)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "Project Roadmap: Option Period 2 & 3"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = theme["title_text"]

    # Table
    total_rows = 2 + len(TASKS)
    total_cols = 1 + MONTH_COLS
    table_shape = slide.shapes.add_table(
        total_rows, total_cols,
        Inches(TABLE_LEFT), Inches(TABLE_TOP),
        Inches(PHASE_COL_WIDTH + CHART_WIDTH),
        Inches(HEADER_HEIGHT * 2 + len(TASKS) * ROW_HEIGHT),
    )
    table = table_shape.table

    # Header rows
    for r in range(2):
        for c in range(total_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["header_bg"]
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = theme["header_text"]
            para.font.size = Pt(10)

    # Row 0: Quarters
    table.cell(0, 0).text = ""
    for i, q in enumerate(QUARTERS):
        c_start = 1 + i * 3
        c_end = 1 + (i + 1) * 3 - 1
        table.cell(0, c_start).text = q
        if c_end > c_start:
            table.cell(0, c_start).merge(table.cell(0, c_end))

    # Row 1: Months
    table.cell(1, 0).text = ""
    for c, m in enumerate(MONTHS):
        table.cell(1, 1 + c).text = m
        para = table.cell(1, 1 + c).text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER

    # Phase labels with vertical merge
    row = 2
    prev_phase = None
    merge_start = None
    for task in TASKS:
        phase_name = task[0]
        cell = table.cell(row, 0)
        cell.text = phase_name if phase_name != prev_phase else ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["phase_bg"]
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = theme["phase_text"]
        para.font.size = Pt(9)
        if phase_name != prev_phase:
            if merge_start is not None:
                table.cell(merge_start, 0).merge(table.cell(row - 1, 0))
            merge_start = row
            prev_phase = phase_name
        row += 1
    if merge_start is not None:
        table.cell(merge_start, 0).merge(table.cell(row - 1, 0))

    # Data cell backgrounds (month columns only)
    for r in range(2, total_rows):
        bg = theme["row_odd"] if (r - 2) % 2 == 1 else theme["row_even"]
        for c in range(1, total_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Task bars
    chart_top = TABLE_TOP + HEADER_HEIGHT * 2
    for i, (phase_name, task_label, start_m, end_m) in enumerate(TASKS):
        row_top = chart_top + i * ROW_HEIGHT
        bar_top = row_top + (ROW_HEIGHT - BAR_HEIGHT) / 2
        bar_left = CHART_LEFT + start_m * MONTH_WIDTH + 0.02
        bar_width = (end_m - start_m + 1) * MONTH_WIDTH - 0.04

        color = bar_colors.get(_phase_key(phase_name), bar_colors["Other"])
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_left), Inches(bar_top),
            Inches(bar_width), Inches(BAR_HEIGHT),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = task_label[:50] + ("..." if len(task_label) > 50 else "")
        p.font.size = Pt(8)
        p.font.color.rgb = theme["bar_text"]
        p.alignment = PP_ALIGN.CENTER

    # Save
    out_dir = Path(__file__).resolve().parent / "output"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / f"DLA_Roadmap_OP2_OP3_{theme_name}.pptx"
    prs.save(str(out_path))
    print(f"Created {out_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate DLA OP2/OP3 roadmap slide")
    parser.add_argument(
        "--theme", choices=["dark", "light"], default="dark",
        help="Color theme (default: dark)",
    )
    args = parser.parse_args()
    main(args.theme)
