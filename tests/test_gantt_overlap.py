"""
Unit tests for Gantt chart overlap detection.

Ensures bars and milestones on Gantt slides never overlap vertically,
regardless of deck content.
"""
import sys
import unittest
from pathlib import Path

# Add project root to path
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation

# MSO_SHAPE_TYPE: 1 = msoAutoShape (bars, diamonds)
_MSO_AUTO_SHAPE = 1

# Minimum gap (inches) between consecutive Gantt bars
# Must be positive to avoid overlap; 0.02 ensures visible separation
MIN_GAP_INCHES = 0.02

# Gantt chart area: bars are right of Deliverable column (chart_left ~2.75)
_GANTT_CHART_LEFT = 2.5


def _get_gantt_bar_shapes(slide):
    """Return bar and diamond shapes from the Gantt chart area."""
    bars = []
    for shape in slide.shapes:
        if not hasattr(shape, "shape_type"):
            continue
        # msoAutoShape = our bars/diamonds
        if shape.shape_type != _MSO_AUTO_SHAPE:
            continue
        if shape.left.inches < _GANTT_CHART_LEFT:
            continue
        bars.append(shape)
    return bars


def _find_gantt_slides(prs):
    """Return slides that contain a Gantt chart (table with Deliverable + bar shapes)."""
    gantt_slides = []
    for slide in prs.slides:
        has_deliverable_table = False
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                try:
                    cell_text = shape.table.cell(0, 0).text or ""
                    if "Deliverable" in cell_text:
                        has_deliverable_table = True
                        break
                except Exception:
                    pass
        bars = _get_gantt_bar_shapes(slide)
        if has_deliverable_table and bars:
            gantt_slides.append((slide, bars))
    return gantt_slides


def _check_no_overlap(bars):
    """
    Check that no two bars overlap vertically.
    Bars are sorted by top; consecutive bars must have gap >= MIN_GAP_INCHES.
    """
    if len(bars) < 2:
        return []
    # Sort by top position
    sorted_bars = sorted(bars, key=lambda s: s.top.inches)
    overlaps = []
    for i in range(len(sorted_bars) - 1):
        a = sorted_bars[i]
        b = sorted_bars[i + 1]
        a_bottom = a.top.inches + a.height.inches
        b_top = b.top.inches
        gap = b_top - a_bottom
        if gap < MIN_GAP_INCHES:
            overlaps.append(
                {
                    "index": i,
                    "a_top": a.top.inches,
                    "a_bottom": a_bottom,
                    "b_top": b_top,
                    "gap": gap,
                    "overlap": a_bottom - b_top,
                }
            )
    return overlaps


class TestGanttOverlap(unittest.TestCase):
    """Test that Gantt charts have no overlapping bars."""

    def setUp(self):
        self.output_dir = ROOT / "output"
        self.output_dir.mkdir(exist_ok=True)

    def test_example_deck_no_overlap(self):
        """Example deck Gantt slide has no overlapping bars."""
        from generate_deck import main, DECK

        main("dark", DECK)
        path = self.output_dir / "Example_Deck_dark.pptx"
        self.assertTrue(path.exists(), f"Expected {path} to exist")
        prs = Presentation(str(path))
        gantt_slides = _find_gantt_slides(prs)
        self.assertGreater(len(gantt_slides), 0, "Expected at least one Gantt slide")
        for slide, bars in gantt_slides:
            overlaps = _check_no_overlap(bars)
            self.assertEqual(
                overlaps,
                [],
                f"Gantt bars overlap: {overlaps}",
            )

    def test_stress_deck_no_overlap(self):
        """Stress deck Gantt slide has no overlapping bars."""
        from generate_deck import main, STRESS_DECK

        main("dark", STRESS_DECK)
        path = self.output_dir / "Stress_Test_dark.pptx"
        self.assertTrue(path.exists(), f"Expected {path} to exist")
        prs = Presentation(str(path))
        gantt_slides = _find_gantt_slides(prs)
        self.assertGreater(len(gantt_slides), 0, "Expected at least one Gantt slide")
        for slide, bars in gantt_slides:
            overlaps = _check_no_overlap(bars)
            self.assertEqual(
                overlaps,
                [],
                f"Gantt bars overlap: {overlaps}",
            )

    def test_light_deck_no_overlap(self):
        """Light theme Gantt slide has no overlapping bars."""
        from generate_deck import main, DECK

        main("light", DECK)
        path = self.output_dir / "Example_Deck_light.pptx"
        self.assertTrue(path.exists(), f"Expected {path} to exist")
        prs = Presentation(str(path))
        gantt_slides = _find_gantt_slides(prs)
        self.assertGreater(len(gantt_slides), 0, "Expected at least one Gantt slide")
        for slide, bars in gantt_slides:
            overlaps = _check_no_overlap(bars)
            self.assertEqual(
                overlaps,
                [],
                f"Gantt bars overlap: {overlaps}",
            )


if __name__ == "__main__":
    unittest.main()
