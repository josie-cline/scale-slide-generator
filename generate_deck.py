#!/usr/bin/env python3
"""
Generate a styled PowerPoint deck for any Scale AI presentation.

Supports multiple slide layouts — title, section, content, two-column,
metrics, table, and gantt — all rendered in Scale styling with dark
and light themes.

Usage:
    python3 generate_deck.py                # dark theme (default)
    python3 generate_deck.py --theme light

Customization:
    Edit the DECK definition below, or ask Cursor:

        "Build me a deck with a title slide, 3 content slides, and a
         metrics slide. Here's the content: ..."
"""
import argparse
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ═══════════════════════════════════════════════════════════════════════════
# DECK DEFINITION — Edit this section for your presentation
# ═══════════════════════════════════════════════════════════════════════════

DECK = {
    "filename": "Example_Deck",
    "slides": [
        # ── Title slide ──────────────────────────────────────────────
        {
            "layout": "title",
            "title": "Quarterly Business Review",
            "subtitle": "Scale AI · Q1 2026",
        },

        # ── Section divider ──────────────────────────────────────────
        {
            "layout": "section",
            "title": "Executive Summary",
        },

        # ── Content slide (title + bullets) ──────────────────────────
        {
            "layout": "content",
            "title": "Key Highlights",
            "bullets": [
                "Deployed platform to production environment",
                "Completed 6 of 6 planned workflow applications",
                "User adoption reached 200+ accounts",
                "Achieved 99.97% platform uptime for the quarter",
            ],
        },

        # ── Two-column comparison ────────────────────────────────────
        {
            "layout": "two_column",
            "title": "Impact: Before & After",
            "left_title": "Manual Process",
            "left_bullets": [
                "3+ hours per report",
                "Inconsistent formatting",
                "Error-prone data entry",
                "Single output per cycle",
            ],
            "right_title": "With Scale Tooling",
            "right_bullets": [
                "Under 30 seconds",
                "Consistent Scale styling",
                "Data-driven generation",
                "Dark + light themes instantly",
            ],
        },

        # ── Metrics / KPI cards ──────────────────────────────────────
        {
            "layout": "metrics",
            "title": "Key Performance Indicators",
            "metrics": [
                {"label": "Active Users", "value": "200+", "detail": "+45% QoQ"},
                {"label": "Applications", "value": "9", "detail": "6 in v1 production"},
                {"label": "Uptime", "value": "99.97%", "detail": "SLA target: 99.5%"},
                {"label": "Avg Response", "value": "< 2s", "detail": "P95 latency"},
            ],
        },

        # ── Table slide ──────────────────────────────────────────────
        {
            "layout": "table",
            "title": "Deliverables Tracker",
            "headers": ["Deliverable", "Owner", "Status", "Due Date"],
            "rows": [
                ["Platform deployment", "Engineering", "Complete", "15 Jan 2026"],
                ["Data integration", "Data Team", "In Progress", "15 Apr 2026"],
                ["User training", "Field Ops", "Planned", "1 May 2026"],
                ["Documentation package", "PM", "In Progress", "30 Apr 2026"],
            ],
        },

        # ── Gantt / roadmap slide ────────────────────────────────────
        {
            "layout": "gantt",
            "title": "Project Roadmap",
            "subtitle": "Scale AI · Option Period 2 & 3",
            "quarters": ["Q1 2026", "Q2 2026", "Q3 2026"],
            "months": ["Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"],
            "phases": [
                "Development",
                "Integration",
                "Sustainment",
                "Deliverables",
            ],
            "tasks": [
                ("Development", "Application refinement", 0, 1, False, "1 Mar 2026"),
                ("Integration", "Data source connection", 0, 3, False, "TBD"),
                ("Integration", "New data connections", 4, 6, False, None),
                ("Sustainment", "Platform maintenance", 0, 6, False, None),
                ("Deliverables", "Kickoff", 0, 0, True, "20 Feb 2026"),
                ("Deliverables", "OP2 Closeout", 2, 3, False, "14 May 2026"),
                ("Deliverables", "OP3 Closeout", 5, 6, False, "14 Aug 2026"),
                ("Deliverables", "Monthly reviews", 0, 6, False, None),
            ],
        },
    ],
}


# ═══════════════════════════════════════════════════════════════════════════
# THEMES — Scale brand colors extracted from .pptx templates
# ═══════════════════════════════════════════════════════════════════════════

FONT_HEADING = "Host Grotesk"
FONT_BODY = "Host Grotesk Light"
FONT_NUMBER = "Manrope"

_ACCENT_PALETTE = [
    RGBColor(0xC8, 0x8B, 0xC4),
    RGBColor(0x7B, 0x8F, 0xDD),
    RGBColor(0x86, 0xBF, 0xF2),
    RGBColor(0xD1, 0xAA, 0xD7),
    RGBColor(0xBB, 0xDE, 0xF2),
    RGBColor(0xA7, 0x7A, 0xFF),
    RGBColor(0xC1, 0x5F, 0xB1),
]

THEMES = {
    "dark": {
        "slide_bg":       RGBColor(0x00, 0x00, 0x00),
        "title_text":     RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text":  RGBColor(0xD9, 0xD9, 0xD9),
        "body_text":      RGBColor(0xDA, 0xDA, 0xDA),
        "muted_text":     RGBColor(0x66, 0x66, 0x66),
        "header_bg":      RGBColor(0x30, 0x30, 0x30),
        "header_text":    RGBColor(0xFF, 0xFF, 0xFF),
        "row_even":       RGBColor(0x00, 0x00, 0x00),
        "row_odd":        RGBColor(0x30, 0x30, 0x30),
        "card_bg":        RGBColor(0x30, 0x30, 0x30),
        "card_border":    RGBColor(0x5E, 0x5E, 0x5E),
        "accent":         RGBColor(0xC8, 0x8B, 0xC4),
        "divider":        RGBColor(0x5E, 0x5E, 0x5E),
        "bullet_color":   RGBColor(0xC8, 0x8B, 0xC4),
        "accent_palette": _ACCENT_PALETTE,
    },
    "light": {
        "slide_bg":       RGBColor(0xFF, 0xFF, 0xFF),
        "title_text":     RGBColor(0x00, 0x00, 0x00),
        "subtitle_text":  RGBColor(0x66, 0x66, 0x66),
        "body_text":      RGBColor(0x66, 0x66, 0x66),
        "muted_text":     RGBColor(0x99, 0x99, 0x99),
        "header_bg":      RGBColor(0xF3, 0xF3, 0xF3),
        "header_text":    RGBColor(0x00, 0x00, 0x00),
        "row_even":       RGBColor(0xFF, 0xFF, 0xFF),
        "row_odd":        RGBColor(0xF3, 0xF3, 0xF3),
        "card_bg":        RGBColor(0xF3, 0xF3, 0xF3),
        "card_border":    RGBColor(0xD9, 0xD9, 0xD9),
        "accent":         RGBColor(0xC8, 0x8B, 0xC4),
        "divider":        RGBColor(0xD9, 0xD9, 0xD9),
        "bullet_color":   RGBColor(0xC8, 0x8B, 0xC4),
        "accent_palette": _ACCENT_PALETTE,
    },
}


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE RENDERERS — You shouldn't need to edit below this line
# ═══════════════════════════════════════════════════════════════════════════

W = 10
H = 5.625    # 16:9 widescreen
MARGIN = 0.5
CONTENT_LEFT = 0.85
SLIDE_BOTTOM = H - 0.25

_ASSETS = Path(__file__).resolve().parent / "assets"
_BG_DARK = _ASSETS / "bg_dark.png"
_BG_LIGHT = _ASSETS / "bg_light.png"
_LOGO_DARK = _ASSETS / "scale_logo_dark.png"
_LOGO_LIGHT = _ASSETS / "scale_logo_light.png"


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_wallpaper(slide, theme_name):
    img_path = _BG_DARK if theme_name == "dark" else _BG_LIGHT
    if not img_path.exists():
        return
    try:
        size_kb = img_path.stat().st_size / 1024
        if size_kb > 500:
            from PIL import Image
            with Image.open(img_path) as im:
                im = im.copy()
                im.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                if im.mode in ("RGBA", "P"):
                    bg = Image.new("RGB", im.size, (255, 255, 255))
                    if im.mode == "P":
                        im = im.convert("RGBA")
                    bg.paste(im, mask=im.split()[-1] if im.mode == "RGBA" else None)
                    im = bg
                elif im.mode != "RGB":
                    im = im.convert("RGB")
                tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
                im.save(tmp.name, "JPEG", quality=88, optimize=True)
                path_to_use = tmp.name
        else:
            path_to_use = str(img_path.resolve())
    except Exception:
        path_to_use = str(img_path.resolve())
    pic = slide.shapes.add_picture(path_to_use, Inches(0), Inches(0), Inches(W), Inches(H))
    spTree = slide.shapes._spTree
    pic_elem = pic._element
    spTree.remove(pic_elem)
    spTree.insert(2, pic_elem)


def _add_branding(slide, theme, theme_name):
    _add_text(slide, W - 2.8, H - 0.35, 2.6, 0.25,
              "Confidential  |  \u00a92025 Scale Inc.",
              size=9, color=theme["muted_text"], font_name=FONT_NUMBER,
              align=PP_ALIGN.RIGHT)
    logo_path = _LOGO_DARK if theme is THEMES["dark"] else _LOGO_LIGHT
    if not logo_path.exists():
        return
    try:
        from PIL import Image
        with Image.open(logo_path) as im:
            w, h = im.size
            ar = w / h
            logo_w = min(0.35, 0.35 * ar)
            logo_h = logo_w / ar
    except Exception:
        logo_w = logo_h = 0.35
    slide.shapes.add_picture(str(logo_path.resolve()),
                             Inches(0.26), Inches(0.32),
                             Inches(logo_w), Inches(logo_h))


def _add_text(slide, left, top, width, height, text, *,
              size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
              font_name=None, word_wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name or (FONT_HEADING if bold else FONT_BODY)
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return box


def _fit_title(text, max_size, width_in, min_size=14):
    cpi = 72 / (max_size * 0.6)
    max_chars = int(cpi * width_in)
    if len(text) <= max_chars:
        return max_size
    return max(min_size, int(max_size * max_chars / len(text)))


def _render_title(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 28, text_w)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT), Inches(2.0), Inches(1.2), Inches(0.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()

    title_h = max(0.6, title_sz / 72 * 1.4)
    _add_text(slide, CONTENT_LEFT, 2.2, text_w, title_h,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])
    if data.get("subtitle"):
        sub_top = min(2.2 + title_h + 0.05, SLIDE_BOTTOM - 0.4)
        _add_text(slide, CONTENT_LEFT, sub_top, text_w, 0.4,
                  data["subtitle"], size=14, color=theme["subtitle_text"])


def _render_section(slide, data, theme):
    text_w = W - 2 * MARGIN
    title_sz = _fit_title(data["title"], 28, text_w)

    _add_text(slide, MARGIN, 2.0, text_w, 0.8,
              data["title"], size=title_sz, bold=True,
              color=theme["title_text"], align=PP_ALIGN.CENTER)

    bar_w = 1.5
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches((W - bar_w) / 2), Inches(2.8), Inches(bar_w), Inches(0.04),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()


def _render_content(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 20, text_w)
    _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT), Inches(0.85), Inches(text_w), Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["divider"]
    line.line.fill.background()

    bullets = data.get("bullets", [])
    n = len(bullets)
    if n == 0:
        return

    avail = SLIDE_BOTTOM - 1.1
    spacing = max(0.28, min(0.50, avail / n))
    font_sz = max(9, min(14, int(14 * spacing / 0.50)))
    dot_sz = max(0.06, 0.1 * (spacing / 0.50))

    top = 1.1
    for bullet in bullets[:int(avail / spacing)]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(CONTENT_LEFT + 0.05), Inches(top + (spacing - dot_sz) / 2),
            Inches(dot_sz), Inches(dot_sz),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme["bullet_color"]
        dot.line.fill.background()

        _add_text(slide, CONTENT_LEFT + 0.3, top, text_w - 0.3, spacing,
                  bullet, size=font_sz, color=theme["body_text"])
        top += spacing


def _render_two_column(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 20, text_w)
    _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT), Inches(0.85), Inches(text_w), Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["divider"]
    line.line.fill.background()

    left_bullets = data.get("left_bullets", [])
    right_bullets = data.get("right_bullets", [])
    max_n = max(len(left_bullets), len(right_bullets), 1)

    col_w = (text_w - 0.5) / 2
    avail = SLIDE_BOTTOM - 1.6
    spacing = max(0.24, min(0.42, avail / max_n))
    font_sz = max(8, min(12, int(12 * spacing / 0.42)))
    dot_sz = max(0.05, 0.09 * (font_sz / 12))

    for col_idx, (title_key, col_bullets) in enumerate([
        ("left_title", left_bullets),
        ("right_title", right_bullets),
    ]):
        x = CONTENT_LEFT + col_idx * (col_w + 0.5)
        _add_text(slide, x, 1.1, col_w, 0.4,
                  data.get(title_key, ""), size=15, bold=True,
                  color=theme["accent"])

        top = 1.6
        for bullet in col_bullets[:int(avail / spacing)]:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.05), Inches(top + (spacing - dot_sz) / 2),
                Inches(dot_sz), Inches(dot_sz),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = theme["bullet_color"]
            dot.line.fill.background()

            _add_text(slide, x + 0.28, top, col_w - 0.28, spacing,
                      bullet, size=font_sz, color=theme["body_text"])
            top += spacing


def _render_metrics(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 20, text_w)
    _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT), Inches(0.85), Inches(text_w), Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["divider"]
    line.line.fill.background()

    metrics = data.get("metrics", [])
    count = len(metrics)
    if count == 0:
        return

    palette = theme["accent_palette"]
    gap = 0.25 if count <= 4 else max(0.1, 0.25 * (4 / count))
    total_gap = gap * (count - 1)
    card_w = max(0.8, (text_w - total_gap) / count)
    card_h = min(1.8, SLIDE_BOTTOM - 1.2)
    card_top = 1.2

    for i, m in enumerate(metrics):
        x = CONTENT_LEFT + i * (card_w + gap)
        inner_pad = min(0.15, card_w * 0.08)
        inner_w = card_w - 2 * inner_pad

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(card_top), Inches(card_w), Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = theme["card_bg"]
        card.line.color.rgb = theme["card_border"]
        card.line.width = Pt(1)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.1), Inches(card_top + card_h / 2 - 0.025),
            Inches(card_w - 0.2), Inches(0.05),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = palette[i % len(palette)]
        accent.line.fill.background()

        val_sz = _fit_title(m["value"], 24, inner_w, min_size=14)
        _add_text(slide, x + inner_pad, card_top + card_h * 0.12, inner_w, card_h * 0.35,
                  m["value"], size=val_sz, bold=True,
                  color=theme["accent"], align=PP_ALIGN.CENTER)

        _add_text(slide, x + inner_pad, card_top + card_h * 0.50, inner_w, card_h * 0.22,
                  m["label"], size=11, bold=True,
                  color=theme["subtitle_text"], align=PP_ALIGN.CENTER)

        if m.get("detail"):
            _add_text(slide, x + inner_pad, card_top + card_h * 0.72, inner_w, card_h * 0.20,
                      m["detail"], size=10,
                      color=theme["muted_text"], align=PP_ALIGN.CENTER)


def _render_table(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 20, text_w)
    _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    if not headers:
        return

    n_cols = len(headers)
    n_rows = 1 + len(rows)
    tbl_w = text_w
    avail = SLIDE_BOTTOM - 1.0
    row_h = max(0.25, min(0.4, avail / n_rows))

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(CONTENT_LEFT), Inches(1.0),
        Inches(tbl_w), Inches(row_h * n_rows),
    )
    table = table_shape.table

    max_lens = [len(str(h)) for h in headers]
    for rd in rows:
        for c, val in enumerate(rd):
            max_lens[c] = max(max_lens[c], len(str(val)))
    total_len = max(sum(max_lens), 1)
    col_widths = [max(0.6, tbl_w * (ml / total_len)) for ml in max_lens]
    ws = sum(col_widths)
    col_widths = [cw * tbl_w / ws for cw in col_widths]
    for c in range(n_cols):
        table.columns[c].width = Inches(col_widths[c])

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["header_bg"]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = FONT_HEADING
        p.font.color.rgb = theme["header_text"]

    for r, row_data in enumerate(rows):
        bg = theme["row_odd"] if r % 2 == 1 else theme["row_even"]
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = FONT_BODY
            p.font.color.rgb = theme["body_text"]


def _strip_table_borders(table):
    """Remove table style and all cell borders for clean overlay rendering."""
    from pptx.oxml.ns import qn
    tbl_pr = table._tbl.tblPr
    for attr in ("bandRow", "bandCol", "firstRow", "lastRow", "firstCol", "lastCol"):
        tbl_pr.set(attr, "0")
    style_id = tbl_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tblStyleId")
    if style_id is not None:
        tbl_pr.remove(style_id)

    for row in table.rows:
        for cell in row.cells:
            tc_pr = cell._tc.get_or_add_tcPr()
            for side in ("lnL", "lnR", "lnT", "lnB"):
                el = tc_pr.find(qn(f"a:{side}"))
                if el is not None:
                    tc_pr.remove(el)
                ln = tc_pr.makeelement(qn(f"a:{side}"), {})
                ln.append(ln.makeelement(qn("a:noFill"), {}))
                tc_pr.append(ln)


def _render_gantt(slide, data, theme):
    months = data.get("months", [])
    quarters = data.get("quarters", [])
    phases = data.get("phases", [])
    tasks = data.get("tasks", [])

    palette = theme["accent_palette"]
    bar_colors = {ph: palette[i % len(palette)] for i, ph in enumerate(phases)}
    fallback = RGBColor(0x99, 0x99, 0x99)

    tbl_left = CONTENT_LEFT
    due_w = 0.9
    n_months = len(months) or 1
    longest_task = max((len(t[1]) for t in tasks), default=10)
    phase_w = max(1.9, min(2.8, longest_task * 0.075))

    max_table_w = W - tbl_left - MARGIN
    month_w = max(0.35, (max_table_w - phase_w - due_w) / n_months)
    chart_w = month_w * n_months
    chart_left = tbl_left + phase_w
    title_w = phase_w + chart_w + due_w

    if data.get("subtitle"):
        _add_text(slide, tbl_left, 0.12, title_w, 0.20,
                  data["subtitle"], size=9, color=theme["subtitle_text"])

    gantt_title_sz = _fit_title(data["title"], 20, title_w)
    _add_text(slide, tbl_left, 0.38, title_w, 0.42,
              data["title"], size=gantt_title_sz, bold=True, color=theme["title_text"])

    tbl_top = 0.95
    n_tasks = len(tasks)
    hdr_h = 0.28
    avail = SLIDE_BOTTOM - tbl_top - hdr_h * 2 - 0.05
    row_h = max(0.28, min(0.32, avail / max(n_tasks, 1)))
    max_tasks = int(avail / row_h)
    if n_tasks > max_tasks:
        tasks = tasks[:max_tasks]
        n_tasks = len(tasks)

    bar_h = min(row_h * 0.75, row_h - 0.08)
    bar_h = max(0.08, bar_h)
    bar_pad = 0.03

    total_rows = 2 + n_tasks
    total_cols = 1 + len(months) + 1
    due_col = total_cols - 1

    table_shape = slide.shapes.add_table(
        total_rows, total_cols,
        Inches(tbl_left), Inches(tbl_top),
        Inches(phase_w + chart_w + due_w),
        Inches(hdr_h * 2 + n_tasks * row_h),
    )
    table = table_shape.table
    _strip_table_borders(table)

    # Set text FIRST, then format (prevents formatting from being wiped)
    table.cell(0, 0).text = "Deliverable"
    if quarters and months:
        nq = len(quarters)
        for i, q in enumerate(quarters):
            s = 1 + i * len(months) // nq
            e = 1 + (i + 1) * len(months) // nq - 1
            if e >= s:
                table.cell(0, s).text = q
                if e > s:
                    table.cell(0, s).merge(table.cell(0, e))
    table.cell(0, due_col).text = "Due Date"
    table.cell(0, due_col).merge(table.cell(1, due_col))
    table.cell(1, 0).text = ""
    for c, m in enumerate(months):
        table.cell(1, 1 + c).text = m[:3] if month_w >= 0.55 else m[:2]

    # Now format headers
    scale = min(1.0, row_h / 0.32)
    hdr_font = max(7, int(10 * scale))
    cell_font = max(6, int(9 * scale))
    for r in range(2):
        for c in range(total_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["header_bg"]
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = theme["header_text"]
            p.font.size = Pt(hdr_font)
            p.font.name = FONT_HEADING
            p.font.bold = (r == 0)
            p.alignment = PP_ALIGN.CENTER

    table.columns[0].width = Inches(phase_w)
    for c in range(1, due_col):
        table.columns[c].width = Inches(month_w)
    table.columns[due_col].width = Inches(due_w)
    table.rows[0].height = Inches(hdr_h)
    table.rows[1].height = Inches(hdr_h)
    for r in range(2, total_rows):
        table.rows[r].height = Inches(row_h)

    task_font = _fit_title("x" * max((len(t[1]) for t in tasks), default=10),
                           cell_font, phase_w - 0.1, min_size=6)
    for row in range(2, total_rows):
        task = tasks[row - 2]
        bg = theme["row_odd"] if (row - 2) % 2 == 1 else theme["row_even"]
        for c in range(total_cols):
            cell = table.cell(row, c)
            if c == 0:
                cell.text = task[1]
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
            elif c == due_col and task[5]:
                cell.text = task[5]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            if c == 0:
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(task_font)
                p.font.name = FONT_BODY
                p.font.color.rgb = theme["body_text"]
                p.alignment = PP_ALIGN.LEFT
            elif c == due_col:
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(cell_font)
                p.font.name = FONT_BODY
                p.font.color.rgb = theme["muted_text"]
                p.alignment = PP_ALIGN.RIGHT

    # Derive bar positions from actual stored EMU values in the table so that
    # float→EMU rounding in the layout math can never misalign bars with rows.
    chart_left_emu  = Inches(chart_left)
    chart_right_emu = Inches(chart_left + chart_w)
    month_w_emu     = Inches(month_w)
    bar_pad_emu     = Inches(bar_pad)
    chart_top_emu   = table_shape.top + table.rows[0].height + table.rows[1].height
    row_y_emu       = chart_top_emu

    for i, task in enumerate(tasks):
        phase, _, start_m, end_m, is_milestone, _ = task
        color = bar_colors.get(phase, fallback)
        row_h_emu   = table.rows[2 + i].height
        row_top_emu = row_y_emu

        if is_milestone:
            sz_emu = min(Inches(0.16), row_h_emu - Inches(0.08))
            sz_emu = max(sz_emu, Inches(0.04))
            cx_emu = chart_left_emu + int((start_m + 0.5) * month_w_emu)
            cx_emu = min(max(cx_emu, chart_left_emu + sz_emu // 2),
                         chart_right_emu - sz_emu // 2)
            cy_emu = row_top_emu + row_h_emu // 2
            shape = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                cx_emu - sz_emu // 2, cy_emu - sz_emu // 2, sz_emu, sz_emu,
            )
        else:
            bar_h_emu = min(int(row_h_emu * 0.75), row_h_emu - Inches(0.08))
            bar_h_emu = max(bar_h_emu, Inches(0.08))
            bt_emu = row_top_emu + (row_h_emu - bar_h_emu) // 2
            bt_emu = max(bt_emu, chart_top_emu + Inches(0.01))
            bl_emu = chart_left_emu + int(start_m * month_w_emu) + bar_pad_emu
            bw_emu = max(Inches(0.2),
                         int((end_m - start_m + 1) * month_w_emu) - 2 * bar_pad_emu)
            if bl_emu + bw_emu > chart_right_emu - bar_pad_emu:
                bw_emu = max(Inches(0.2), chart_right_emu - bar_pad_emu - bl_emu)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bl_emu, bt_emu, bw_emu, bar_h_emu,
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        row_y_emu += row_h_emu


# ── Charts (bar, pie, line) via python-pptx native chart API ─────────────

def _render_chart(slide, data, theme):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 20, text_w)
    _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT), Inches(0.85), Inches(text_w), Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["divider"]
    line.line.fill.background()

    chart_type_str = data.get("chart_type", "bar")
    chart_types = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "donut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_type = chart_types.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])
    for series in data.get("series", []):
        chart_data.add_series(series["name"], series["values"])

    chart_x = Inches(CONTENT_LEFT)
    chart_y = Inches(1.1)
    chart_cx = Inches(text_w)
    chart_cy = Inches(SLIDE_BOTTOM - 1.1 - 0.1)

    chart_frame = slide.shapes.add_chart(
        xl_type, chart_x, chart_y, chart_cx, chart_cy, chart_data,
    )
    chart = chart_frame.chart

    chart.has_legend = len(data.get("series", [])) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = theme["body_text"]

    palette = theme["accent_palette"]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = palette[i % len(palette)]

    is_pie = chart_type_str in ("pie", "donut")

    if not is_pie:
        try:
            cat_axis = chart.category_axis
            cat_axis.tick_labels.font.size = Pt(9)
            cat_axis.tick_labels.font.color.rgb = theme["body_text"]
            cat_axis.tick_labels.font.name = FONT_BODY
            cat_axis.format.line.fill.background()
        except (ValueError, AttributeError):
            pass

        try:
            val_axis = chart.value_axis
            val_axis.tick_labels.font.size = Pt(8)
            val_axis.tick_labels.font.color.rgb = theme["muted_text"]
            val_axis.tick_labels.font.name = FONT_NUMBER
            val_axis.format.line.fill.background()
            val_axis.major_gridlines.format.line.color.rgb = theme["divider"]
            val_axis.major_gridlines.format.line.width = Pt(0.5)
        except (ValueError, AttributeError):
            pass

    plot = chart.plots[0]
    if not is_pie:
        plot.gap_width = 80
    else:
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(10)
        plot.data_labels.font.color.rgb = theme["title_text"]
        plot.data_labels.number_format = '0%'
        plot.data_labels.show_percentage = True
        plot.data_labels.show_value = False


# ── Quote slide ──────────────────────────────────────────────────────────

def _render_quote(slide, data, theme):
    quote_text = data.get("quote", "")
    text_w = W - 2 * MARGIN - 1.5
    quote_x = MARGIN + 0.75

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN + 0.35), Inches(1.8), Inches(0.05), Inches(2.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()

    quote_sz = _fit_title(quote_text, 24, text_w, min_size=14)
    _add_text(slide, quote_x, 1.8, text_w, 2.0,
              f"\u201c{quote_text}\u201d", size=quote_sz,
              color=theme["title_text"], align=PP_ALIGN.LEFT, italic=True)

    if data.get("attribution"):
        _add_text(slide, quote_x, 4.0, text_w, 0.35,
                  f"\u2014 {data['attribution']}", size=13,
                  color=theme["accent"])


# ── Fact / stat callout ──────────────────────────────────────────────────

def _render_fact(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN

    if data.get("title"):
        title_sz = _fit_title(data["title"], 20, text_w)
        _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
                  data["title"], size=title_sz, bold=True, color=theme["title_text"])
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(CONTENT_LEFT), Inches(0.85), Inches(text_w), Inches(0.015),
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = theme["divider"]
        slide.shapes[-1].line.fill.background()

    if data.get("label"):
        _add_text(slide, CONTENT_LEFT, 1.4, text_w, 0.35,
                  data["label"], size=14, bold=True,
                  color=theme["subtitle_text"], align=PP_ALIGN.CENTER)

    value_sz = _fit_title(data.get("value", ""), 56, text_w - 1.0, min_size=28)
    _add_text(slide, CONTENT_LEFT, 1.9, text_w, 1.0,
              data.get("value", ""), size=value_sz, bold=True,
              color=theme["accent"], font_name=FONT_NUMBER, align=PP_ALIGN.CENTER)

    bar_w = 1.5
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT + (text_w - bar_w) / 2), Inches(3.15),
        Inches(bar_w), Inches(0.04),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = theme["accent"]
    slide.shapes[-1].line.fill.background()

    if data.get("detail"):
        _add_text(slide, CONTENT_LEFT, 3.4, text_w, 0.45,
                  data["detail"], size=14,
                  color=theme["body_text"], align=PP_ALIGN.CENTER)


# ── Process / flow ───────────────────────────────────────────────────────

def _render_process(slide, data, theme):
    text_w = W - CONTENT_LEFT - MARGIN
    title_sz = _fit_title(data["title"], 20, text_w)
    _add_text(slide, CONTENT_LEFT, 0.3, text_w, 0.5,
              data["title"], size=title_sz, bold=True, color=theme["title_text"])

    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CONTENT_LEFT), Inches(0.85), Inches(text_w), Inches(0.015),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = theme["divider"]
    slide.shapes[-1].line.fill.background()

    steps = data.get("steps", [])
    n = len(steps)
    if n == 0:
        return

    palette = theme["accent_palette"]
    gap = 0.12
    step_w = (text_w - gap * (n - 1)) / n
    step_w = min(step_w, 2.8)
    step_h = 1.0
    y_top = 1.4

    for i, step in enumerate(steps):
        x = CONTENT_LEFT + i * (step_w + gap)
        color = palette[i % len(palette)]

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y_top), Inches(step_w), Inches(step_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        _add_text(slide, x + 0.2, y_top + 0.15, step_w - 0.4, 0.30,
                  str(i + 1).zfill(2), size=24, bold=True,
                  color=theme["slide_bg"])

        label_sz = _fit_title(step.get("label", ""), 14, step_w - 0.4, min_size=10)
        _add_text(slide, x + 0.2, y_top + 0.50, step_w - 0.4, 0.40,
                  step.get("label", ""), size=label_sz, bold=True,
                  color=theme["slide_bg"])

        if step.get("detail"):
            _add_text(slide, x + 0.05, y_top + step_h + 0.25, step_w - 0.1, 0.45,
                      step["detail"], size=min(11, label_sz),
                      color=theme["body_text"], align=PP_ALIGN.CENTER)


RENDERERS = {
    "title": _render_title,
    "section": _render_section,
    "content": _render_content,
    "two_column": _render_two_column,
    "metrics": _render_metrics,
    "table": _render_table,
    "gantt": _render_gantt,
    "chart": _render_chart,
    "quote": _render_quote,
    "fact": _render_fact,
    "process": _render_process,
}


def main(theme_name: str = "dark", deck: dict = None):
    deck = deck or DECK
    theme = THEMES[theme_name]

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    for slide_data in deck["slides"]:
        layout = slide_data.get("layout", "content")
        renderer = RENDERERS.get(layout)
        if not renderer:
            print(f"Warning: unknown layout '{layout}', skipping")
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, theme["slide_bg"])
        _add_wallpaper(slide, theme_name)
        renderer(slide, slide_data, theme)
        _add_branding(slide, theme, theme_name)
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    out_dir = Path(__file__).resolve().parent / "output"
    out_dir.mkdir(exist_ok=True)
    filename = deck.get("filename", "Deck")
    out_path = out_dir / f"{filename}_{theme_name}.pptx"
    prs.save(str(out_path))
    print(f"Created {out_path}")
    return str(out_path)


STRESS_DECK = {
    "filename": "Stress_Test",
    "slides": [
        {"layout": "title", "title": "Stress Test: Long Title That Should Shrink", "subtitle": "Overflow validation"},
        {"layout": "content", "title": "15 Bullets", "bullets": [f"Bullet point number {i+1} with extra text" for i in range(15)]},
        {"layout": "two_column", "title": "Dense Two Column", "left_title": "Left", "left_bullets": [f"Left {i+1}" for i in range(12)], "right_title": "Right", "right_bullets": [f"Right {i+1}" for i in range(12)]},
        {"layout": "metrics", "title": "8 KPI Cards", "metrics": [{"label": f"Metric {i+1}", "value": f"{(i+1)*111}", "detail": f"Detail {i+1}"} for i in range(8)]},
        {"layout": "table", "title": "Large Table", "headers": ["Col A", "Col B", "Col C", "Col D"], "rows": [[f"R{r}C{c}" for c in range(4)] for r in range(15)]},
        {"layout": "gantt", "title": "Dense Gantt", "subtitle": "Stress test", "quarters": ["Q1", "Q2", "Q3", "Q4"], "months": ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"], "phases": ["Phase A", "Phase B", "Phase C"], "tasks": [("Phase A", f"Task {i+1}", i % 12, min(i % 12 + 2, 11), False, None) for i in range(15)]},
    ],
}


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a styled Scale AI slide deck")
    parser.add_argument("--theme", choices=["dark", "light"], default="dark")
    parser.add_argument("--stress", action="store_true")
    args = parser.parse_args()
    deck = STRESS_DECK if args.stress else DECK
    main(args.theme, deck=deck)
